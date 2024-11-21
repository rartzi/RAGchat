# utils/text_utils.py
import tiktoken
from typing import List
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
import streamlit as st

def count_tokens(text: str) -> int:
    """Count tokens in a text string."""
    encoding = tiktoken.encoding_for_model("gpt-4")
    return len(encoding.encode(text))

def chunk_text(text: str, max_tokens: int = 500) -> List[str]:
    """Split text into chunks of specified maximum tokens."""
    chunks = []
    current_chunk = []
    current_size = 0
    
    sentences = text.replace('\n', ' ').split('.')
    
    for sentence in sentences:
        sentence = sentence.strip() + '.'
        sentence_tokens = count_tokens(sentence)
        
        if sentence_tokens > max_tokens:
            # Handle long sentences
            words = sentence.split()
            current_word_chunk = []
            current_word_size = 0
            
            for word in words:
                word_size = count_tokens(word + ' ')
                if current_word_size + word_size > max_tokens:
                    chunks.append(' '.join(current_word_chunk))
                    current_word_chunk = [word]
                    current_word_size = word_size
                else:
                    current_word_chunk.append(word)
                    current_word_size += word_size
            
            if current_word_chunk:
                chunks.append(' '.join(current_word_chunk))
                
        elif current_size + sentence_tokens > max_tokens:
            chunks.append(' '.join(current_chunk))
            current_chunk = [sentence]
            current_size = sentence_tokens
        else:
            current_chunk.append(sentence)
            current_size += sentence_tokens
    
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    
    return chunks

def read_document(file_path: str) -> str:
    """Extract text from various document formats."""
    try:
        if file_path.endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
                
        elif file_path.endswith('.pdf'):
            reader = PdfReader(file_path)
            return '\n'.join(page.extract_text() for page in reader.pages)
            
        elif file_path.endswith('.docx'):
            doc = docx.Document(file_path)
            return '\n'.join(paragraph.text for paragraph in doc.paragraphs)
            
        elif file_path.endswith('.pptx'):
            presentation = Presentation(file_path)
            return '\n'.join(
                shape.text for slide in presentation.slides 
                for shape in slide.shapes if hasattr(shape, "text")
            )
            
        return ""
    except Exception as e:
        st.error(f"Error reading {file_path}: {str(e)}")
        return ""