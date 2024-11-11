# 1_database_manager.py

import os
import streamlit as st
from openai import OpenAI
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
import chromadb
import uuid
import tiktoken
from typing import List
import requests
from datetime import datetime

# Configuration
DB_DIR = "/Users/ronenartzi/development/playground/gui/streamlit/RAGchat/db"
DEFAULT_DOCS_DIR = "/Users/ronenartzi/development/playground/gui/streamlit/RAGchat/documents"

# Initialize ChromaDB
os.makedirs(DB_DIR, exist_ok=True)
chroma_client = chromadb.PersistentClient(path=DB_DIR)
collection = chroma_client.get_or_create_collection(
    name="document_embeddings",
    metadata={"hnsw:space": "cosine"}
)

def count_tokens(text: str) -> int:
    """Count tokens in a text string."""
    encoding = tiktoken.encoding_for_model("gpt-4")
    return len(encoding.encode(text))

def chunk_text(text: str, max_tokens: int = 500) -> List[str]:
    """Split text into chunks of specified maximum tokens."""
    chunks = []
    current_chunk = []
    current_size = 0
    
    # Split text into sentences
    sentences = text.replace('\n', ' ').split('.')
    
    for sentence in sentences:
        sentence = sentence.strip() + '.'
        sentence_tokens = count_tokens(sentence)
        
        if sentence_tokens > max_tokens:
            # If single sentence is too long, split it into smaller chunks
            words = sentence.split()
            current_word_chunk = []
            current_word_size = 0
            
            for word in words:
                word_size = count_tokens(word + ' ')
                if current_word_size + word_size > max_tokens:
                    chunks.append(' '.join(current_word_chunk))
                    current_word_chunk = [word]
                    current_word_size = word_size
                else:
                    current_word_chunk.append(word)
                    current_word_size += word_size
            
            if current_word_chunk:
                chunks.append(' '.join(current_word_chunk))
                
        elif current_size + sentence_tokens > max_tokens:
            chunks.append(' '.join(current_chunk))
            current_chunk = [sentence]
            current_size = sentence_tokens
        else:
            current_chunk.append(sentence)
            current_size += sentence_tokens
    
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    
    return chunks

def read_document(file_path: str) -> str:
    """Extract text from various document formats."""
    try:
        if file_path.endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
                
        elif file_path.endswith('.pdf'):
            reader = PdfReader(file_path)
            text = ''
            for page in reader.pages:
                text += page.extract_text() + '\n'
            return text
            
        elif file_path.endswith('.docx'):
            doc = docx.Document(file_path)
            return '\n'.join(paragraph.text for paragraph in doc.paragraphs)
            
        elif file_path.endswith('.pptx'):
            presentation = Presentation(file_path)
            text = ''
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'
            return text
            
        return ""
    except Exception as e:
        st.error(f"Error reading {file_path}: {str(e)}")
        return ""

def get_embeddings(text: str, provider: str = "openai") -> List[float]:
    """Generate embeddings using specified provider."""
    try:
        if provider == "openai":
            client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
            if not client.api_key:
                raise ValueError("OpenAI API key not found!")
                
            response = client.embeddings.create(
                input=text,
                model="text-embedding-ada-002"
            )
            return response.data[0].embedding
            
        else:  # ollama
            response = requests.post(
                "http://localhost:11434/api/embeddings",
                json={"model": "nomic-embed-text", "prompt": text}
            )
            if response.status_code != 200:
                raise ValueError(f"Ollama error: {response.status_code}")
            return response.json()['embedding']
            
    except Exception as e:
        st.error(f"Error generating embeddings: {str(e)}")
        return None

def get_database_stats():
    """Get current database statistics."""
    try:
        results = collection.get()
        if not results or not results['ids']:
            return {
                "document_count": 0,
                "chunk_count": 0,
                "sources": [],
                "status": "empty"
            }
        
        # Count only actual documents (exclude any special entries)
        valid_documents = []
        for meta in results['metadatas']:
            if meta and isinstance(meta, dict) and 'source' in meta:
                valid_documents.append(meta['source'])
                
        sources = set(valid_documents)
        return {
            "document_count": len(sources),
            "chunk_count": len(valid_documents),
            "sources": sorted(list(sources)),
            "status": "ready"
        }
    except Exception as e:
        return {
            "document_count": 0,
            "chunk_count": 0,
            "sources": [],
            "status": f"error: {str(e)}"
        }

def process_and_add_document(file_path: str, provider: str) -> bool:
    """Process and add a document to the database."""
    try:
        content = read_document(file_path)
        if not content.strip():
            st.warning(f"No content extracted from {file_path}")
            return False
            
        # Get chunks
        chunks = chunk_text(content)
        file_name = os.path.basename(file_path)
        
        success_count = 0
        for i, chunk in enumerate(chunks):
            embedding = get_embeddings(chunk, provider)
            if embedding:
                chunk_id = f"{file_name}_chunk_{i}_{str(uuid.uuid4())}"
                collection.add(
                    ids=[chunk_id],
                    embeddings=[embedding],
                    metadatas=[{
                        "source": file_name,
                        "chunk": i,
                        "chunk_total": len(chunks),
                        "content": chunk,
                        "added_at": datetime.now().isoformat()
                    }],
                    documents=[chunk]
                )
                success_count += 1
            else:
                st.warning(f"Failed to generate embedding for chunk {i} of {file_name}")
        
        if success_count > 0:
            st.success(f"Added {success_count} chunks from {file_name}")
            return True
        return False
        
    except Exception as e:
        st.error(f"Error processing {file_path}: {str(e)}")
        return False

# UI Setup
st.set_page_config(
    page_title="Document Database Manager",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS
st.markdown("""
    <style>
        .stAlert {
            margin-top: 10px;
            margin-bottom: 10px;
        }
        .stProgress {
            margin-top: 20px;
            margin-bottom: 20px;
        }
        .stButton {
            margin-top: 15px;
            margin-bottom: 15px;
        }
        .file-list {
            margin-top: 10px;
            margin-bottom: 10px;
            padding: 10px;
            background-color: #f0f2f6;
            border-radius: 5px;
        }
    </style>
""", unsafe_allow_html=True)

st.title("📚 Document Database Manager")

# Sidebar
with st.sidebar:
    st.header("Settings")
    provider = st.radio(
        "Embedding Provider",
        ["OpenAI", "Ollama"],
        help="Select the embedding provider to use",
        key="provider"
    )
    
    # Show token usage warning for OpenAI
    if provider == "OpenAI":
        st.info("⚠️ Note: Using OpenAI will incur token usage costs")
    
    st.divider()
    stats = get_database_stats()
    st.header("Database Stats")
    
    if stats["status"] == "empty":
        st.info("Database is empty")
    elif stats["status"] == "ready":
        st.metric("📑 Documents", stats["document_count"])
        st.metric("🔢 Total Chunks", stats["chunk_count"])
    else:
        st.error(f"Database Error: {stats['status']}")

# Main content
col1, col2 = st.columns([2, 1])

with col1:
    st.header("Add Documents")
    docs_dir = st.text_input(
        "Documents Directory:",
        value=DEFAULT_DOCS_DIR,
        help="Enter the path to your documents directory"
    )
    
    if docs_dir and os.path.isdir(docs_dir):
        files = [f for f in os.listdir(docs_dir) 
                if f.endswith(('.txt', '.pdf', '.docx', '.pptx'))]
        
        if files:
            st.success(f"Found {len(files)} documents")
            with st.expander("View Files", expanded=True):
                for file in files:
                    st.text(f"📄 {file}")
                
            if st.button("🔄 Process Documents", type="primary", use_container_width=True):
                progress = st.progress(0)
                processed_files = 0
                
                for i, file in enumerate(files):
                    with st.spinner(f"Processing {file}..."):
                        status = process_and_add_document(
                            os.path.join(docs_dir, file), 
                            provider.lower()
                        )
                        if status:
                            processed_files += 1
                    progress.progress((i + 1) / len(files))
                
                if processed_files > 0:
                    st.success(f"✅ Successfully processed {processed_files} out of {len(files)} documents!")
                    st.balloons()
                else:
                    st.error("❌ No documents were successfully processed")
                st.rerun()
        else:
            st.warning("No supported documents found in directory")
    else:
        st.error("Invalid directory path")

with col2:
    st.header("Manage Database")
    
    if st.button("🗑️ Clear Database", type="secondary", use_container_width=True):
        try:
            all_ids = collection.get()['ids']
            if all_ids:
                collection.delete(ids=all_ids)
                st.success("✨ Database cleared!")
                st.rerun()
            else:
                st.info("Database is already empty")
        except Exception as e:
            st.error(f"Error clearing database: {str(e)}")
    
    if stats["document_count"] > 0:
        st.subheader("Current Documents")
        for source in stats["sources"]:
            st.text(f"📄 {source}")

# Footer
st.divider()
st.caption(f"Database Location: {DB_DIR}")

# Instructions
with st.expander("ℹ️ Instructions", expanded=False):
    st.markdown("""
    ### How to use:
    1. **Select Embedding Provider** in the sidebar
        - OpenAI (requires API key)
        - Ollama (requires local installation)
    
    2. **Prepare Your Documents**
        - Place documents in the specified directory
        - Supported formats: PDF, DOCX, PPTX, TXT
    
    3. **Process Documents**
        - Click "Process Documents" to add them to the database
        - Wait for processing to complete
    
    4. **Manage Database**
        - View current documents in the database
        - Clear database if needed
        
    ### Supported File Types:
    - PDF files (*.pdf)
    - Word documents (*.docx)
    - PowerPoint presentations (*.pptx)
    - Text files (*.txt)
    """)